"""Document Generator Agent — creates professional PDF CVs and PPTX materials."""

import os
import structlog
from typing import Optional
from app.config import settings

logger = structlog.get_logger()


def _x(text) -> str:
    """Escape text for safe use inside ReportLab Paragraph (XML + encoding safe)."""
    from xml.sax.saxutils import escape
    if text is None:
        return ""
    s = str(text)
    # Encode to latin-1 (Windows-1252 superset) — the encoding used by ReportLab's
    # standard Helvetica/Times fonts. Replace unsupported chars with '?' to avoid
    # KeyError / UnicodeEncodeError inside doc.build().
    s = s.encode("latin-1", errors="replace").decode("latin-1")
    return escape(s)


def _build_cv_elements(tailored_data: dict) -> list:
    """Build ReportLab flowable elements from tailored CV data."""
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, Spacer, HRFlowable
    from reportlab.lib.colors import HexColor

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('CVTitle', parent=styles['Title'],
                                 fontSize=18, textColor=HexColor('#1a1a2e'), spaceAfter=6)
    heading_style = ParagraphStyle('CVHeading', parent=styles['Heading2'],
                                   fontSize=13, textColor=HexColor('#6366F1'),
                                   spaceAfter=4, spaceBefore=12)
    body_style = ParagraphStyle('CVBody', parent=styles['Normal'],
                                fontSize=10, spaceAfter=3, leading=14)
    bullet_style = ParagraphStyle('CVBullet', parent=body_style,
                                  bulletIndent=15, leftIndent=25)

    cv_data = tailored_data.get("tailored_cv", tailored_data) or {}
    elements = []

    # Name
    personal = cv_data.get("personal_info", {}) or {}
    elements.append(Paragraph(_x(personal.get("name") or "Candidate"), title_style))

    # Contact info
    contact_parts = []
    for field in ["email", "phone", "location", "linkedin", "github", "portfolio"]:
        val = personal.get(field)
        if val:
            contact_parts.append(_x(val))
    for field in ["date_of_birth", "nationality", "gender", "marital_status", "visa_status"]:
        val = personal.get(field)
        if val:
            contact_parts.append(f"{field.replace('_', ' ').title()}: {_x(val)}")
    if contact_parts:
        elements.append(Paragraph(" | ".join(contact_parts), body_style))

    elements.append(Spacer(1, 10))
    elements.append(HRFlowable(width="100%", color=HexColor('#6366F1')))

    # Summary
    summary = cv_data.get("summary", "")
    if summary:
        elements.append(Paragraph("PROFESSIONAL SUMMARY", heading_style))
        elements.append(Paragraph(_x(summary), body_style))

    # Skills
    skills = cv_data.get("skills", {})
    if skills:
        elements.append(Paragraph("SKILLS", heading_style))
        if isinstance(skills, dict):
            for category, skill_list in skills.items():
                if category.startswith("_") or category == "all":
                    continue
                if isinstance(skill_list, list) and skill_list:
                    safe_skills = ", ".join(_x(s) for s in skill_list if s)
                    if safe_skills:
                        elements.append(Paragraph(
                            f"<b>{_x(category.title())}:</b> {safe_skills}", body_style))
        elif isinstance(skills, list) and skills:
            safe_skills = ", ".join(_x(s) for s in skills if s)
            if safe_skills:
                elements.append(Paragraph(safe_skills, body_style))

    # Experience
    experience = cv_data.get("experience") or []
    if experience:
        elements.append(Paragraph("EXPERIENCE", heading_style))
        for exp in experience:
            if not isinstance(exp, dict):
                continue
            elements.append(Paragraph(
                f"<b>{_x(exp.get('role', ''))}</b> — {_x(exp.get('company', ''))} ({_x(exp.get('duration', ''))})",
                body_style))
            for ach in (exp.get("achievements") or []):
                if ach and isinstance(ach, str):
                    elements.append(Paragraph(f"• {_x(ach)}", bullet_style))

    # Education
    education = cv_data.get("education") or []
    if education:
        elements.append(Paragraph("EDUCATION", heading_style))
        for edu in education:
            if not isinstance(edu, dict):
                continue
            elements.append(Paragraph(
                f"<b>{_x(edu.get('degree', ''))}</b> — {_x(edu.get('institution', ''))} ({_x(edu.get('year', ''))})",
                body_style))

    # Projects
    projects = cv_data.get("projects") or []
    if projects:
        elements.append(Paragraph("PROJECTS", heading_style))
        for proj in projects:
            if not isinstance(proj, dict):
                continue
            techs = ", ".join(_x(t) for t in (proj.get("technologies") or []) if t)
            elements.append(Paragraph(
                f"<b>{_x(proj.get('name', ''))}</b> — {_x(proj.get('description', ''))}",
                body_style))
            if techs:
                elements.append(Paragraph(f"Technologies: {techs}", bullet_style))

    # Certifications
    certifications = cv_data.get("certifications") or []
    if certifications:
        elements.append(Paragraph("CERTIFICATIONS", heading_style))
        for cert in certifications:
            if isinstance(cert, dict) and cert.get("name"):
                name = _x(cert.get("name", ""))
                issuer = _x(cert.get("issuer") or "")
                date = _x(cert.get("issue_date") or "")
                line = f"<b>{name}</b>"
                if issuer:
                    line += f" — {issuer}"
                if date:
                    line += f" ({date})"
                elements.append(Paragraph(f"• {line}", bullet_style))
            elif isinstance(cert, str) and cert:
                elements.append(Paragraph(f"• {_x(cert)}", bullet_style))

    # Languages
    languages = cv_data.get("languages") or []
    if languages:
        elements.append(Paragraph("LANGUAGES", heading_style))
        lang_parts = []
        for lang in languages:
            if isinstance(lang, dict) and lang.get("language"):
                prof = lang.get("proficiency")
                lang_str = _x(lang["language"])
                if prof:
                    lang_str += f" ({_x(prof)})"
                lang_parts.append(lang_str)
            elif isinstance(lang, str) and lang:
                lang_parts.append(_x(lang))
        if lang_parts:
            elements.append(Paragraph(", ".join(lang_parts), body_style))

    return elements


async def generate_cv_pdf_bytes(tailored_data: dict) -> Optional[bytes]:
    """Generate a CV PDF and return raw bytes (no disk I/O).

    This is the primary generation path used when attaching to emails.
    """
    try:
        import io
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate

        elements = _build_cv_elements(tailored_data)
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4,
                                leftMargin=0.75*inch, rightMargin=0.75*inch,
                                topMargin=0.5*inch, bottomMargin=0.5*inch)
        doc.build(elements)
        pdf_bytes = buffer.getvalue()
        logger.info("pdf_bytes_generated", size=len(pdf_bytes))
        return pdf_bytes

    except Exception as e:
        logger.error("pdf_bytes_generation_error", error=str(e), exc_info=True)
        # Hard fallback: minimal canvas PDF — never fails on encoding or XML issues
        try:
            import io
            from reportlab.pdfgen import canvas as rl_canvas
            from reportlab.lib.pagesizes import A4
            cv_data = (tailored_data or {}).get("tailored_cv", tailored_data) or {}
            personal = (cv_data.get("personal_info") or {})
            name = str(personal.get("name") or "Candidate").encode("latin-1", errors="replace").decode("latin-1")
            buf = io.BytesIO()
            c = rl_canvas.Canvas(buf, pagesize=A4)
            w, h = A4
            c.setFont("Helvetica-Bold", 18)
            c.drawString(50, h - 70, name[:80])
            c.setFont("Helvetica", 11)
            c.drawString(50, h - 100, "Tailored CV — Digital FTE")
            c.save()
            logger.info("pdf_fallback_generated")
            return buf.getvalue()
        except Exception as e2:
            logger.error("pdf_fallback_failed", error=str(e2), exc_info=True)
            return None


async def generate_cv_pdf(tailored_data: dict, template: str = "modern") -> Optional[str]:
    """Generate a CV PDF, write to disk, and return the file path."""
    pdf_bytes = await generate_cv_pdf_bytes(tailored_data)
    if not pdf_bytes:
        return None
    try:
        import uuid
        generated_dir = os.path.abspath(settings.GENERATED_DIR)
        os.makedirs(generated_dir, exist_ok=True)
        output_path = os.path.join(generated_dir, f"cv_{uuid.uuid4().hex[:8]}.pdf")
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        logger.info("pdf_generated", path=output_path)
        return output_path
    except Exception as e:
        logger.error("pdf_write_error", error=str(e), exc_info=True)
        return None


async def generate_interview_pptx(prep_data: dict, job_title: str, company: str) -> Optional[str]:
    """Generate PPTX study slides from interview prep data."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import uuid

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"Interview Prep: {job_title}"
        slide.placeholders[1].text = f"Company: {company}"

        # Company Research
        research = prep_data.get("company_research", {})
        if research:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Company Research"
            body = slide.placeholders[1]
            for key, val in research.items():
                if val:
                    p = body.text_frame.add_paragraph()
                    p.text = f"{key.replace('_', ' ').title()}: {val}"
                    p.font.size = Pt(14)

        # Technical Questions
        tech_qs = prep_data.get("technical_questions", [])
        for i, q in enumerate(tech_qs[:5]):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Technical Q{i+1}"
            body = slide.placeholders[1]
            body.text = f"Q: {q.get('question', '')}\n\nA: {q.get('answer', '')}"

        os.makedirs(settings.GENERATED_DIR, exist_ok=True)
        output_path = os.path.join(settings.GENERATED_DIR, f"prep_{uuid.uuid4().hex[:8]}.pptx")
        prs.save(output_path)
        logger.info("pptx_generated", path=output_path)
        return output_path

    except Exception as e:
        logger.error("pptx_generation_error", error=str(e))
        return None
